"""Generate DTDL Talent Hack presentation for AI Experiment Copilot."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Brand palette
NAVY = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
ACCENT_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

OUTPUT = Path(__file__).resolve().parents[1] / "docs" / "DTDL_Talent_Hack_Presentation.pptx"


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=Inches(0), height=Inches(0.08)) -> None:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0),
        top,
        Inches(13.333),
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12), Inches(0.4))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def style_title(text_frame, size=44, color=DARK_TEXT, bold=True, align=PP_ALIGN.LEFT):
    text_frame.word_wrap = True
    for i, para in enumerate(text_frame.paragraphs):
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(size if i == 0 else size - 8)
            run.font.color.rgb = color
            run.font.bold = bold


def add_bullets(text_frame, items, size=18, color=DARK_TEXT, spacing=Pt(8)):
    text_frame.clear()
    text_frame.word_wrap = True
    for idx, item in enumerate(items):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.line_spacing = 1.25


def build_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide, top=Inches(3.55), height=Inches(0.06))

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.4))
    tf = title_box.text_frame
    tf.text = "DTDL Talent Hack"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.85), Inches(11.5), Inches(1.2))
    stf = sub_box.text_frame
    stf.text = "AI Experiment Copilot & Decision Intelligence"
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.font.size = Pt(24)
    sp.font.color.rgb = ACCENT_LIGHT

    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.6))
    ttf = tag_box.text_frame
    ttf.text = "Universal A/B Testing & Decision Platform"
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tp.font.size = Pt(16)
    tp.font.color.rgb = MUTED


def build_problem_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide)

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.7))
    title.text_frame.text = "Problem Statement"
    style_title(title.text_frame, size=36, color=NAVY)

    # Context box
    ctx = slide.shapes.add_shape(1, Inches(0.7), Inches(1.35), Inches(12), Inches(1.05))
    ctx.fill.solid()
    ctx.fill.fore_color.rgb = ACCENT_LIGHT
    ctx.line.color.rgb = ACCENT
    ctx_tf = ctx.text_frame
    ctx_tf.word_wrap = True
    ctx_tf.margin_left = Inches(0.2)
    ctx_tf.margin_right = Inches(0.2)
    ctx_tf.margin_top = Inches(0.12)
    ctx_p = ctx_tf.paragraphs[0]
    ctx_p.text = (
        "A/B experiments require manual setup, statistical expertise, and constant monitoring — "
        "leading to inconsistent quality and slow decisions."
    )
    ctx_p.font.size = Pt(16)
    ctx_p.font.color.rgb = DARK_TEXT
    ctx_p.line_spacing = 1.3

    # Challenge
    ch = slide.shapes.add_textbox(Inches(0.7), Inches(2.65), Inches(12), Inches(0.5))
    ch.text_frame.text = "How might we help?"
    ch.text_frame.paragraphs[0].font.size = Pt(20)
    ch.text_frame.paragraphs[0].font.bold = True
    ch.text_frame.paragraphs[0].font.color.rgb = ACCENT

    bullets = slide.shapes.add_textbox(Inches(0.7), Inches(3.15), Inches(5.8), Inches(3.5))
    add_bullets(
        bullets.text_frame,
        [
            "Guide teams through the full experiment lifecycle",
            "Generate hypotheses from business goals",
            "Recommend metrics, audiences & configurations",
            "Validate setup before launch",
            "Monitor performance in real time",
        ],
        size=16,
    )

    bullets2 = slide.shapes.add_textbox(Inches(6.8), Inches(3.15), Inches(5.8), Inches(3.5))
    add_bullets(
        bullets2.text_frame,
        [
            "Explain why variants win or underperform",
            "Summarize results in business-friendly language",
            "Recommend next actions: Scale / Continue / Stop / Rollback",
            "Reduce manual effort across the lifecycle",
            "Deliver explainable, data-driven decisions",
        ],
        size=16,
    )

    # Success criteria footer
    foot = slide.shapes.add_textbox(Inches(0.7), Inches(6.55), Inches(12), Inches(0.55))
    foot_tf = foot.text_frame
    foot_tf.text = (
        "Success: faster experiment creation · high-quality configs · accurate insights · "
        "reliable recommendations · measurable eval metrics (creation time, acceptance rate, analysis time)"
    )
    foot_tf.paragraphs[0].font.size = Pt(11)
    foot_tf.paragraphs[0].font.color.rgb = SLATE
    foot_tf.paragraphs[0].line_spacing = 1.2

    add_footer(slide, "AI Experiment Copilot")


def build_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide)

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.7))
    title.text_frame.text = "Solution Architecture"
    style_title(title.text_frame, size=36, color=NAVY)

    tagline = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(12), Inches(0.4))
    tagline.text_frame.text = (
        "Six microservices · shared Postgres · LangGraph agent with Playwright + SQL tools"
    )
    tagline.text_frame.paragraphs[0].font.size = Pt(13)
    tagline.text_frame.paragraphs[0].font.color.rgb = SLATE

    # Two-column service layout
    left_col = [
        ("PM Dashboard (:5174)", "Chat, hypothesis generation, config recommendations, one-click analyze, apply verdict"),
        ("Copilot Backend (:3001)", "FastAPI · LangGraph ReAct agent · lifecycle & eval APIs"),
        ("Playwright MCP (:8931)", "Opens variant URLs · visual diff · metric inference from page context"),
    ]
    right_col = [
        ("E-Commerce Storefront (:5173)", "Mock test subject · A/B variant rendering · funnel events"),
        ("E-Com Backend (:3002)", "Event ingestion · feature flags · migrations & seed data"),
        ("PostgreSQL (:5432)", "universal_events table · experiments · read-only agent SQL role"),
    ]

    col_w = Inches(5.85)
    box_h = Inches(1.05)
    gap = Inches(0.18)
    start_y = Inches(1.65)

    for col_idx, items in enumerate([left_col, right_col]):
        x = Inches(0.7) if col_idx == 0 else Inches(6.75)
        for row_idx, (name, desc) in enumerate(items):
            top = start_y + row_idx * (box_h + gap)
            box = slide.shapes.add_shape(1, x, top, col_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = ACCENT if col_idx == 0 else NAVY
            box.line.fill.background()
            tf = box.text_frame
            tf.margin_left = Inches(0.18)
            tf.margin_top = Inches(0.1)
            tf.word_wrap = True
            p1 = tf.paragraphs[0]
            p1.text = name
            p1.font.size = Pt(14)
            p1.font.bold = True
            p1.font.color.rgb = WHITE
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(11)
            p2.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            p2.space_before = Pt(2)

    # Agent workflow strip
    wf_box = slide.shapes.add_shape(1, Inches(0.7), Inches(5.55), Inches(11.9), Inches(0.95))
    wf_box.fill.solid()
    wf_box.fill.fore_color.rgb = ACCENT_LIGHT
    wf_box.line.color.rgb = ACCENT
    wf_tf = wf_box.text_frame
    wf_tf.margin_left = Inches(0.2)
    wf_tf.margin_top = Inches(0.12)
    wf_tf.word_wrap = True
    wf_p1 = wf_tf.paragraphs[0]
    wf_p1.text = "Agent Analysis Pipeline"
    wf_p1.font.size = Pt(13)
    wf_p1.font.bold = True
    wf_p1.font.color.rgb = NAVY
    wf_p2 = wf_tf.add_paragraph()
    wf_p2.text = (
        "Inspect variants → Discover events (SQL) → Infer success metric → "
        "Two-proportion z-test → Scale / Continue / Stop / Rollback"
    )
    wf_p2.font.size = Pt(12)
    wf_p2.font.color.rgb = DARK_TEXT
    wf_p2.space_before = Pt(4)

    # Key design point
    kp = slide.shapes.add_textbox(Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.35))
    kp.text_frame.text = (
        "Key design: metric is inferred (not pre-configured) · LLM writes SQL & prose · "
        "statistics & verdicts are deterministic · DB-layer read-only guardrails"
    )
    kp.text_frame.paragraphs[0].font.size = Pt(10)
    kp.text_frame.paragraphs[0].font.color.rgb = SLATE

    add_footer(slide, "AI Experiment Copilot")


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title_slide(prs)
    build_problem_slide(prs)
    build_architecture_slide(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
